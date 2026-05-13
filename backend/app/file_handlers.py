from __future__ import annotations

import csv
import re
import shutil
import uuid
import zipfile
from pathlib import Path

import fitz
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from .scanner import scan_text


BASE_DIR = Path(__file__).resolve().parents[1]
STORAGE_DIR = BASE_DIR / "storage"
UPLOAD_DIR = STORAGE_DIR / "uploads"
REDACTED_DIR = STORAGE_DIR / "redacted"

for directory in (UPLOAD_DIR, REDACTED_DIR):
    directory.mkdir(parents=True, exist_ok=True)


SUPPORTED_EXTENSIONS = {".txt", ".csv", ".docx", ".xlsx", ".pdf", ".pptx"}
TEXT_EXTENSIONS = {".txt", ".csv"}
OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
MAX_UPLOAD_BYTES = 20 * 1024 * 1024
MAX_OFFICE_UNCOMPRESSED_BYTES = 100 * 1024 * 1024
MAX_OFFICE_FILE_COUNT = 1000
MAX_EXTRACTED_TEXT_CHARS = 500_000
MAX_PDF_PAGES = 200
FORMULA_PREFIXES = ("=", "+", "-", "@")


def save_upload(filename: str, content: bytes) -> Path:
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的文件类型：{suffix or '未知'}")
    if len(content) > MAX_UPLOAD_BYTES:
        raise ValueError("上传文件过大，最大支持 20MB")
    _validate_file_signature(suffix, content)
    safe_name = f"{uuid.uuid4().hex}{suffix}"
    path = UPLOAD_DIR / safe_name
    path.write_bytes(content)
    return path


def _preview(text: str) -> str:
    compact = "\n".join(line.rstrip() for line in text.splitlines() if line.strip())
    return compact


def redact_file(source: Path, original_filename: str, mode: str = "mask") -> dict:
    suffix = source.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的文件类型：{suffix}")
    _validate_saved_file(source, suffix)

    target = REDACTED_DIR / _redacted_filename(original_filename, suffix)
    if suffix == ".txt":
        text = _limit_text(source.read_text(encoding="utf-8", errors="ignore"))
        redacted, findings = scan_text(text, mode, "全文")
        target.write_text(redacted, encoding="utf-8")
        return {"path": target, "preview": _preview(redacted), "findings": findings}

    if suffix == ".csv":
        return _redact_csv(source, target, mode)
    if suffix == ".xlsx":
        return _redact_xlsx(source, target, mode)
    if suffix == ".docx":
        return _redact_docx(source, target, mode)
    if suffix == ".pdf":
        return _redact_pdf(source, target, mode)
    if suffix == ".pptx":
        return _redact_pptx(source, target, mode)

    raise ValueError(f"不支持的文件类型：{suffix}")


def _redacted_filename(original_filename: str, fallback_suffix: str) -> str:
    original_path = Path(original_filename)
    suffix = original_path.suffix.lower() or fallback_suffix
    stem = original_path.stem or "redacted"
    safe_stem = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", stem).strip(" .")
    if not safe_stem:
        safe_stem = "redacted"
    return f"{safe_stem}-fixed{suffix}"


def _validate_file_signature(suffix: str, content: bytes) -> None:
    if suffix == ".pdf" and not content.startswith(b"%PDF-"):
        raise ValueError("文件扩展名与 PDF 文件内容不匹配")
    if suffix in OFFICE_EXTENSIONS and not content.startswith(b"PK\x03\x04"):
        raise ValueError("文件扩展名与 Office 文件内容不匹配")
    if suffix in TEXT_EXTENSIONS and b"\x00" in content[:4096]:
        raise ValueError("文本文件内容异常，疑似二进制文件")


def _validate_saved_file(path: Path, suffix: str) -> None:
    if suffix in OFFICE_EXTENSIONS:
        _validate_office_zip(path)


def _validate_office_zip(path: Path) -> None:
    try:
        with zipfile.ZipFile(path) as archive:
            infos = archive.infolist()
            if len(infos) > MAX_OFFICE_FILE_COUNT:
                raise ValueError("Office 文件内部文件数量异常，已拒绝处理")
            total_size = 0
            for item in infos:
                name = item.filename.replace("\\", "/")
                if name.startswith("/") or ".." in Path(name).parts:
                    raise ValueError("Office 文件包含异常路径，已拒绝处理")
                total_size += item.file_size
                if total_size > MAX_OFFICE_UNCOMPRESSED_BYTES:
                    raise ValueError("Office 文件解压后体积过大，已拒绝处理")
    except zipfile.BadZipFile as exc:
        raise ValueError("Office 文件结构损坏或格式不正确") from exc


def _limit_text(text: str) -> str:
    if len(text) > MAX_EXTRACTED_TEXT_CHARS:
        raise ValueError("提取文本过长，已拒绝处理")
    return text


def _neutralize_spreadsheet_formula(value: str) -> str:
    stripped = value.lstrip()
    if stripped and stripped[0] in FORMULA_PREFIXES:
        return "'" + value
    return value


def _redact_csv(source: Path, target: Path, mode: str) -> dict:
    findings: list[dict] = []
    preview_rows: list[str] = []
    with source.open("r", encoding="utf-8-sig", errors="ignore", newline="") as src:
        rows = list(csv.reader(src))

    redacted_rows = []
    for row_idx, row in enumerate(rows, start=1):
        new_row = []
        for col_idx, value in enumerate(row, start=1):
            redacted, cell_findings = scan_text(value, mode, f"第 {row_idx} 行，第 {col_idx} 列")
            findings.extend(cell_findings)
            new_row.append(_neutralize_spreadsheet_formula(redacted))
        redacted_rows.append(new_row)
        if row_idx <= 20:
            preview_rows.append(", ".join(new_row))

    with target.open("w", encoding="utf-8-sig", newline="") as dst:
        csv.writer(dst).writerows(redacted_rows)

    return {"path": target, "preview": _preview("\n".join(preview_rows)), "findings": findings}


def _redact_xlsx(source: Path, target: Path, mode: str) -> dict:
    workbook = load_workbook(source)
    findings: list[dict] = []
    preview_lines: list[str] = []

    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            values = []
            for cell in row:
                if isinstance(cell.value, str):
                    redacted, cell_findings = scan_text(cell.value, mode, f"{sheet.title}!{cell.coordinate}")
                    findings.extend(cell_findings)
                    cell.value = _neutralize_spreadsheet_formula(redacted)
                values.append("" if cell.value is None else str(cell.value))
            if len(preview_lines) < 30 and any(values):
                preview_lines.append(f"{sheet.title}: " + " | ".join(values))

    workbook.save(target)
    return {"path": target, "preview": _preview("\n".join(preview_lines)), "findings": findings}


def _redact_docx(source: Path, target: Path, mode: str) -> dict:
    document = Document(source)
    findings: list[dict] = []
    preview_lines: list[str] = []

    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = _limit_text(paragraph.text)
        redacted, para_findings = scan_text(text, mode, f"第 {index} 段")
        if para_findings:
            _replace_paragraph_text(paragraph, redacted)
            findings.extend(para_findings)
        if paragraph.text.strip() and len(preview_lines) < 30:
            preview_lines.append(paragraph.text)

    for table_idx, table in enumerate(document.tables, start=1):
        for row_idx, row in enumerate(table.rows, start=1):
            for col_idx, cell in enumerate(row.cells, start=1):
                text = _limit_text(cell.text)
                redacted, cell_findings = scan_text(text, mode, f"表 {table_idx} 行 {row_idx} 列 {col_idx}")
                if cell_findings:
                    cell.text = redacted
                    findings.extend(cell_findings)

    document.save(target)
    return {"path": target, "preview": _preview("\n".join(preview_lines)), "findings": findings}


def _replace_paragraph_text(paragraph, text: str) -> None:
    for run in paragraph.runs:
        run.text = ""
    if paragraph.runs:
        paragraph.runs[0].text = text
    else:
        paragraph.add_run(text)


def _redact_pdf(source: Path, target: Path, mode: str) -> dict:
    doc = fitz.open(source)
    findings: list[dict] = []
    preview_lines: list[str] = []
    if len(doc) > MAX_PDF_PAGES:
        doc.close()
        raise ValueError("PDF 页数过多，已拒绝处理")

    for page_index in range(len(doc)):
        page = doc[page_index]
        text = _limit_text(page.get_text())
        redacted, page_findings = scan_text(text, mode, f"第 {page_index + 1} 页")
        findings.extend(page_findings)
        if redacted.strip() and len(preview_lines) < 20:
            preview_lines.append(redacted)
        for item in page_findings:
            evidence = item["evidence"]
            for rect in page.search_for(evidence):
                page.add_redact_annot(rect, text=item["replacement"], fill=(1, 1, 1))
        page.apply_redactions()

    doc.save(target)
    doc.close()
    return {"path": target, "preview": _preview("\n".join(preview_lines)), "findings": findings}


def _redact_pptx(source: Path, target: Path, mode: str) -> dict:
    shutil.copy2(source, target)
    presentation = Presentation(target)
    findings: list[dict] = []
    preview_lines: list[str] = []

    for slide_idx, slide in enumerate(presentation.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if not hasattr(shape, "text") or not shape.text:
                continue
            text = _limit_text(shape.text)
            redacted, shape_findings = scan_text(text, mode, f"第 {slide_idx} 页，元素 {shape_idx}")
            if shape_findings:
                shape.text = redacted
                findings.extend(shape_findings)
            if shape.text.strip() and len(preview_lines) < 30:
                preview_lines.append(shape.text)

    presentation.save(target)
    return {"path": target, "preview": _preview("\n".join(preview_lines)), "findings": findings}
